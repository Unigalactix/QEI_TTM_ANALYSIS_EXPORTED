"""
October 2025 TTM Analysis - PowerPoint Generator v2
Creates a professional slide deck matching SlideTemplate style without using images as backgrounds
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Paths
BASE_PATH = r"C:\Users\nigopal\OneDrive - Microsoft\Documents\QEI_TTM_Analysis"
OCT_PATH = os.path.join(BASE_PATH, "OctTTM")
OUTPUT_FILE = os.path.join(OCT_PATH, "October_2025_TTM_Analysis.pptx")

# Create presentation with blank template (16:9 aspect ratio)
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 ratio (1920px wide)
prs.slide_height = Inches(7.5)    # 16:9 ratio (1080px tall)

# Design colors (Microsoft-style palette)
TITLE_COLOR = RGBColor(0, 120, 215)  # Microsoft Blue
ACCENT_COLOR = RGBColor(0, 90, 158)  # Darker Blue
TEXT_COLOR = RGBColor(50, 50, 50)    # Dark Gray
LIGHT_BG = RGBColor(245, 245, 245)   # Light Gray background

def add_title_slide(title, subtitle):
    """Create title slide with professional styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add colored background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add blue accent bar at top
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(0.5)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Subtitle
    left = Inches(1)
    top = Inches(4.2)
    width = Inches(8)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.word_wrap = True
    
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    
    return slide

def add_content_slide(title, content_list):
    """Create content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Blue title bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(1)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content area
    left = Inches(0.75)
    top = Inches(1.5)
    width = Inches(11.8)
    height = Inches(5.5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = item
        p.level = 0
        p.font.size = Pt(16)  # Reduced from 18 to prevent overflow
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)  # Reduced from 12 for tighter spacing
        
        # Add bullet
        p.font.name = 'Calibri'
    
    return slide

def add_image_slide(title, image_path, caption=""):
    """Create slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Blue title bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(1)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add image - calculate size to fit within slide bounds
    if os.path.exists(image_path):
        from PIL import Image
        
        # Get image dimensions
        img = Image.open(image_path)
        img_width, img_height = img.size
        
        # Available space (with margins)
        max_width = Inches(11.5)
        max_height = Inches(5.5) if caption else Inches(6)
        
        # Calculate scaling to fit within bounds while maintaining aspect ratio
        width_scale = max_width / Inches(img_width / 96)  # Assuming 96 DPI
        height_scale = max_height / Inches(img_height / 96)
        scale = min(width_scale, height_scale, 1.0)  # Don't upscale
        
        # Calculate final dimensions
        final_width = Inches(img_width / 96) * scale
        final_height = Inches(img_height / 96) * scale
        
        # Center the image
        left = (Inches(13.333) - final_width) / 2
        top = Inches(1.5)
        
        slide.shapes.add_picture(image_path, left, top, width=final_width, height=final_height)
    
    # Caption if provided
    if caption:
        left = Inches(1)
        top = Inches(6.8)
        width = Inches(11.333)
        height = Inches(0.6)
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        p = caption_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = TEXT_COLOR
    
    return slide

def add_two_column_slide(title, left_content, right_content):
    """Create slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Blue title bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(1)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(5.5)
    
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    # Right column
    left = Inches(6.8)
    right_box = slide.shapes.add_textbox(left, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    return slide

# ============================================================================
# BUILD PRESENTATION
# ============================================================================

print("=" * 80)
print("OCTOBER 2025 TTM ANALYSIS - POWERPOINT GENERATOR V2")
print("=" * 80)
print()

# Slide 1: Title
print("📊 Creating Slide 1: Title Slide...")
add_title_slide(
    "October 2025 TTM Analysis",
    "Quality Engineering Insights - Time to Mitigate Review"
)

# Slide 2: Executive Summary
print("📊 Creating Slide 2: Executive Summary...")
add_content_slide(
    "Executive Summary",
    [
        "📊 Total Incidents: 117 (114 after exclusions)",
        "⏱️ P75 TTM: 190 minutes (baseline for analysis)",
        "🔍 3 Incidents Excluded: BCDR drills and EUAP regions (7,802 minutes)",
        "📈 TTM Impact: Exclusions reduced average TTM by 23% (293→225 min)",
        "🎯 High TTM Multiplier: 17.0x slower than normal incidents (976 vs 58 min)",
        "🔧 Resolution Gap: 46.7% ad-hoc fixes in High TTM vs 20.7% in Normal",
        "💡 What-If: Top 5 event systems account for 36.3% of P75 TTM"
    ]
)

# Slide 3: TTM Distribution
print("📊 Creating Slide 3: TTM Distribution...")
ttm_dist_path = os.path.join(OCT_PATH, "October_TTM_Distribution.png")
add_image_slide(
    "TTM Distribution - October 2025",
    ttm_dist_path,
    "Distribution shows concentration in 0-200 minute range with outliers"
)

# Slide 4: Summary Statistics
print("📊 Creating Slide 4: Summary Statistics...")
add_content_slide(
    "Summary Statistics",
    [
        "Total Incidents: 117",
        "Mean TTM: 292.9 minutes (4.9 hours)",
        "Median (P50): 92.0 minutes",
        "P75: 190.0 minutes",
        "P90: 597.0 minutes",
        "Standard Deviation: 752.8 minutes (high variability)",
        "Range: 2 to 6,095 minutes"
    ]
)

# Slide 5: Top Services
print("📊 Creating Slide 5: Top Services...")
top_services_path = os.path.join(OCT_PATH, "October_Top_Services.png")
add_image_slide(
    "Top Services by TTM - October 2025",
    top_services_path,
    "SQL Control Plane and Xstore dominate TTM minutes"
)

# Slide 6: Daily Timeline
print("📊 Creating Slide 6: Daily Timeline...")
timeline_path = os.path.join(OCT_PATH, "October_Daily_Timeline.png")
add_image_slide(
    "Daily Timeline - October 2025",
    timeline_path,
    "Incident frequency and TTM patterns across the month"
)

# Slide 7: Severity Distribution
print("📊 Creating Slide 7: Severity Distribution...")
severity_path = os.path.join(OCT_PATH, "October_Severity_Distribution.png")
add_image_slide(
    "Severity Distribution - October 2025",
    severity_path,
    "Distribution of incidents by severity level"
)

# Slide 8: Month-over-Month Comparison
print("📊 Creating Slide 8: Month-over-Month Comparison...")
add_two_column_slide(
    "October vs September Comparison",
    [
        "📊 OCTOBER 2025:",
        "• Total Incidents: 117",
        "• P75 TTM: 190.0 min",
        "• Mean: 292.9 min",
        "• Median: 92.0 min",
        "• P90: 597.0 min"
    ],
    [
        "📊 SEPTEMBER 2025:",
        "• Total Incidents: 175",
        "• P75 TTM: 180.0 min",
        "• Mean: 261.3 min",
        "• Median: 88.0 min",
        "• P90: 548.0 min"
    ]
)

# Slide 9: What-If Cumulative Impact
print("📊 Creating Slide 9: What-If Cumulative Impact...")
whatif_cumulative_path = os.path.join(OCT_PATH, "WhatIf_Cumulative_Impact.png")
add_image_slide(
    "What-If Analysis: Cumulative Impact",
    whatif_cumulative_path,
    "P75 TTM reduction as events are removed (Top 5 = 36.3% reduction)"
)

# Slide 10: What-If Marginal Returns
print("📊 Creating Slide 10: What-If Marginal Returns...")
whatif_marginal_path = os.path.join(OCT_PATH, "WhatIf_Cumulative_Marginal.png")
add_image_slide(
    "What-If Analysis: Marginal Returns",
    whatif_marginal_path,
    "Diminishing returns after top 5 events (13.8 min/event → 3.7 min/event)"
)

# Slide 11: What-If Key Findings
print("📊 Creating Slide 11: What-If Key Findings...")
add_content_slide(
    "What-If Analysis: Key Findings",
    [
        "🎯 89 Unique Event Systems: 76 root events + 41 cascading outages",
        "📉 Top 5 Events Impact: 36.3% of P75 TTM (190→121 minutes)",
        "📊 Diminishing Returns: Events 1-5 avg 13.8 min/event, Events 6-10 avg 3.7 min/event",
        "🔍 Single Largest Event: #694602140 (Xstore EUAP, 6,095 min, 32% of P75)",
        "💡 Prevention Priority: Top 10 events = 47.9% of P75 TTM impact",
        "⚠️ Long Tail: 79 events contribute remaining 52.1% (avg 1.3 min/event)"
    ]
)

# Slide 12: Exclusions
print("📊 Creating Slide 12: Exclusions...")
add_content_slide(
    "Exclusions: BCDR and EUAP Incidents",
    [
        "🚫 3 Incidents Excluded (2.6% of total):",
        "",
        "1️⃣ #694602140: Xstore, 6,095 min (EUAP region, 'By Design')",
        "2️⃣ #694752515: Compute RP, 1,372 min (BCDR drill + EUAP)",
        "3️⃣ #694624704: SQL MI, 335 min (EUAP region)",
        "",
        "📊 Impact: 23% reduction in average TTM (293→225 min)",
        "✅ All analysis uses filtered dataset (114 incidents)"
    ]
)

# Slide 13: Narrative Insights
print("📊 Creating Slide 13: Narrative Insights...")
add_content_slide(
    "Narrative Insights: Resolution Gap",
    [
        "🔍 High TTM vs Normal TTM Comparison:",
        "",
        "⏱️ TTM Multiplier: 17.0x slower (976 min vs 58 min)",
        "🔧 Ad-Hoc Resolution: 46.7% (High) vs 20.7% (Normal) = +26.0pp gap",
        "📋 TSG Usage: 60.0% (High) vs 69.0% (Normal) = -9.0pp gap",
        "🤖 BRAIN Detection: 0% in High TTM cohort vs higher in Normal",
        "👥 Human Detection: Dominates High TTM incidents (manual escalation)",
        "💡 Key Insight: Lack of standardized procedures drives extended resolution times"
    ]
)

# Slide 14: Service Patterns
print("📊 Creating Slide 14: Service Patterns...")
add_content_slide(
    "Narrative Insights: Service Patterns",
    [
        "🏆 Top Services by High TTM Impact:",
        "",
        "1️⃣ SQL Control Plane: 12,140 min (41.4% of High TTM total)",
        "2️⃣ Xstore: 7,322 min (25.0% of High TTM total)",
        "3️⃣ Network Infrastructure: 3,845 min (13.1%)",
        "4️⃣ Other Services: 6,013 min (20.5%)",
        "",
        "🔍 Service-Specific Challenges:",
        "• SQL Control Plane: Complex dependencies, regional cascades",
        "• Xstore: Storage layer issues with broad impact"
    ]
)

# Slide 15: Root Cause Patterns
print("📊 Creating Slide 15: Root Cause Patterns...")
add_content_slide(
    "Narrative Insights: Root Cause Patterns",
    [
        "🐛 Software Bugs: 5.9x more common in High TTM (20% vs 3.4%)",
        "⚙️ Configuration Issues: Present in both cohorts, not discriminating",
        "🔄 Deployment Problems: More investigation needed in High TTM",
        "📡 Dependency Failures: External service issues compound resolution time",
        "",
        "🎯 Pattern Analysis:",
        "• High TTM incidents involve novel/complex issues requiring deep investigation",
        "• Standard runbooks and TSGs insufficient for edge cases",
        "• Need enhanced diagnostic tools and knowledge base"
    ]
)

# Slide 16: Recommendations
print("📊 Creating Slide 16: Recommendations...")
add_content_slide(
    "Recommendations: Priority Actions",
    [
        "1️⃣ Create Missing TSGs: Target top 5 event types (36.3% impact potential), focus on SQL Control Plane and Xstore scenarios",
        "",
        "2️⃣ Improve BRAIN Detection: 0% detection in High TTM cohort needs attention, enhance pattern recognition for complex scenarios",
        "",
        "3️⃣ Enhance Code Quality: Address 5.9x higher bug rate in High TTM incidents, strengthen pre-deployment testing for edge cases",
        "",
        "4️⃣ Automate Common Resolutions: Reduce 46.7% ad-hoc resolution rate, build runbook automation for frequent patterns"
    ]
)

# Slide 17: Key Takeaways
print("📊 Creating Slide 17: Key Takeaways...")
add_content_slide(
    "Key Takeaways",
    [
        "✅ October showed 33% fewer incidents than September (117 vs 175)",
        "⚠️ P75 TTM increased slightly (+5.6%, 190 vs 180 min)",
        "🎯 Top 5 event systems represent highest ROI for prevention (36.3% impact)",
        "📉 Diminishing returns after top 10 events (focus prioritization)",
        "🔧 Resolution process gaps (46.7% ad-hoc, 0% BRAIN detection in High TTM)",
        "💡 Service concentration (SQL Control Plane + Xstore = 66% of High TTM)",
        "📋 TSG creation and BRAIN enhancement are critical next steps"
    ]
)

# Slide 18: Appendix
print("📊 Creating Slide 18: Appendix...")
add_title_slide(
    "Appendix",
    "Data Sources, Methodology, and Definitions"
)

# Slide 19: Data Sources
print("📊 Creating Slide 19: Data Sources...")
add_content_slide(
    "Data Sources and Methodology",
    [
        "📊 Data Source:",
        "• Kusto Cluster: icmdataro.centralus.kusto.windows.net",
        "• Database: IcmDataCommon",
        "• Time Period: October 1-31, 2025",
        "• Query: Step 1 of PROMPT_1.md workflow",
        "",
        "🔍 Analysis Methodology:",
        "• 9-Step PROMPT_1.md workflow",
        "• Exclusions applied: BCDR drills and EUAP regions",
        "• Event system analysis using RootResponsibleIncidentId",
        "• Narrative text mining across 13 description fields"
    ]
)

# Slide 20: Definitions
print("📊 Creating Slide 20: Definitions...")
add_content_slide(
    "Key Definitions",
    [
        "⏱️ TTM (Time to Mitigate): Minutes from incident creation to mitigation",
        "📊 P75: 75th percentile (3 of 4 incidents resolve faster)",
        "🎯 Event System: Root incident + all cascading outages (via RootResponsibleIncidentId)",
        "🔴 High TTM: Incidents in Q5 (top 20% by TTM)",
        "🟢 Normal TTM: Incidents in Q1-Q4 (bottom 80%)",
        "🔧 Ad-Hoc Resolution: No documented TSG or runbook used",
        "🤖 BRAIN Detection: Automated anomaly detection system",
        "🚫 Exclusions: BCDR drills and EUAP (pre-production) incidents"
    ]
)

# Save presentation
print()
print(f"💾 Saving presentation to: {OUTPUT_FILE}")
prs.save(OUTPUT_FILE)

print("=" * 80)
print(f"✅ SUCCESS! PowerPoint presentation created with 20 slides")
print(f"📁 Location: {OUTPUT_FILE}")
print("=" * 80)
print()
print("📊 Slide Breakdown:")
print("  1. Title Slide")
print("  2. Executive Summary")
print("  3-7. Visualizations (TTM Distribution, Services, Timeline, Severity, Statistics)")
print("  8. Month-over-Month Comparison")
print("  9-11. What-If Analysis (3 slides)")
print("  12. Exclusions")
print("  13-15. Narrative Insights (3 slides)")
print("  16. Recommendations")
print("  17. Key Takeaways")
print("  18-20. Appendix (3 slides)")
print()
print("=" * 80)
