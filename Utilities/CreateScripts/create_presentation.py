"""
October 2025 TTM Analysis - PowerPoint Generator
Creates a comprehensive slide deck with templates, images, and analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from PIL import Image

# Paths
BASE_PATH = r"C:\Users\nigopal\OneDrive - Microsoft\Documents\QEI_TTM_Analysis"
TEMPLATE_PATH = os.path.join(BASE_PATH, "Utilities", "SlideTemplate")
OCT_PATH = os.path.join(BASE_PATH, "OctTTM")
OUTPUT_FILE = os.path.join(OCT_PATH, "October_2025_TTM_Analysis.pptx")

# Template images
TITLE_TEMPLATE = os.path.join(TEMPLATE_PATH, "Title_Slide.png")
TEXT_TEMPLATE = os.path.join(TEMPLATE_PATH, "Content_Text_Slide.png")
ANALYSIS_TEMPLATE = os.path.join(TEMPLATE_PATH, "Content_Analysis_Slide.png")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_background(slide, template_path):
    """Add template background to slide"""
    left = top = Inches(0)
    slide.shapes.add_picture(template_path, left, top, width=prs.slide_width, height=prs.slide_height)

def add_title_slide(title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, TITLE_TEMPLATE)
    
    # Add title text box
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(60, 60, 60)

def add_content_slide(title, bullet_points):
    """Create content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, TEXT_TEMPLATE)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.level = 0 if not point.startswith("  ") else 1
        p.font.size = Pt(18) if p.level == 0 else Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(10)

def add_image_slide(title, image_path, caption=""):
    """Create slide with image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, ANALYSIS_TEMPLATE)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Image
    if os.path.exists(image_path):
        left = Inches(1.5)
        top = Inches(1.8)
        width = Inches(7)
        slide.shapes.add_picture(image_path, left, top, width=width)
    
    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        p = caption_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 80, 80)

def add_two_column_slide(title, left_content, right_content):
    """Create two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, TEXT_TEMPLATE)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, point in enumerate(left_content):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, point in enumerate(right_content):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(8)

print("="*80)
print("OCTOBER 2025 TTM ANALYSIS - POWERPOINT GENERATOR")
print("="*80)

# Slide 1: Title
print("\n📊 Creating Slide 1: Title Slide...")
add_title_slide(
    "October 2025 TTM Analysis",
    "Time to Mitigate Analysis & Insights"
)

# Slide 2: Executive Summary
print("📊 Creating Slide 2: Executive Summary...")
add_content_slide(
    "Executive Summary",
    [
        "📈 Total Incidents: 117 (filtered: 114 after exclusions)",
        "⏱️ P75 TTM: 190 minutes (3.2 hours)",
        "📉 Average TTM: 293 minutes (4.9 hours)",
        "🎯 Key Finding: High TTM incidents take 17.0x longer to resolve",
        "🚨 Top Impact Services:",
        "  • SQL Control Plane: 12,140 min total TTM",
        "  • Xstore: 7,322 min total TTM",
        "  • Fabric Network Devices: 37 incidents",
        "🔧 Resolution Gap: 46.7% ad-hoc vs 20.7% (High vs Normal TTM)",
        "⚠️ 3 incidents excluded (BCDR drills + EUAP region)"
    ]
)

# Slide 3: TTM Distribution
print("📊 Creating Slide 3: TTM Distribution...")
add_image_slide(
    "TTM Distribution",
    os.path.join(OCT_PATH, "October_TTM_Distribution.png"),
    "P75 = 190 minutes | Mean = 293 minutes | Median = 71 minutes"
)

# Slide 4: Summary Statistics
print("📊 Creating Slide 4: Summary Statistics...")
add_two_column_slide(
    "Summary Statistics",
    [
        "TTM Metrics:",
        "• Mean: 293 min",
        "• Median (P50): 71 min",
        "• P75: 190 min",
        "• P90: 353 min",
        "• Range: 3 - 6,095 min",
        "",
        "Severity Distribution:",
        "• Sev 0: 1 (0.9%)",
        "• Sev 1: 9 (7.7%)",
        "• Sev 2: 107 (91.5%)"
    ],
    [
        "Incident Characteristics:",
        "• Auto-Detection: 0.0%",
        "• Multi-Region: 8 (6.8%)",
        "• Change-Related: 42 (35.9%)",
        "",
        "By Quintile (Q5 = Highest):",
        "• Q1: 12 min avg",
        "• Q2: 40 min avg",
        "• Q3: 70 min avg",
        "• Q4: 152 min avg",
        "• Q5: 1,170 min avg"
    ]
)

# Slide 5: Top Services
print("📊 Creating Slide 5: Top Services...")
add_image_slide(
    "Top Impacted Services",
    os.path.join(OCT_PATH, "October_Top_Services.png"),
    "Fabric Network Devices leads in count (37), SQL Control Plane leads in avg TTM (2,436 min)"
)

# Slide 6: Daily Timeline
print("📊 Creating Slide 6: Daily Timeline...")
add_image_slide(
    "Incident Timeline",
    os.path.join(OCT_PATH, "October_Daily_Timeline.png"),
    "Daily incident distribution throughout October 2025"
)

# Slide 7: Severity Distribution
print("📊 Creating Slide 7: Severity Distribution...")
add_image_slide(
    "Severity Distribution",
    os.path.join(OCT_PATH, "October_Severity_Distribution.png"),
    "Severity 2 dominates at 91.5% of all incidents"
)

# Slide 8: Month-over-Month Comparison
print("📊 Creating Slide 8: Comparison to September...")
add_content_slide(
    "October vs September Comparison",
    [
        "📊 Incident Volume:",
        "  • October: 117 incidents",
        "  • September: 175 incidents",
        "  • Change: -58 incidents (-33.1%)",
        "",
        "⏱️ TTM Metrics:",
        "  • P75 TTM: 190 min (Oct) vs 180 min (Sept) → +10 min (+5.6%)",
        "  • Mean TTM: 293 min (Oct) vs 196 min (Sept) → +97 min (+49.5%)",
        "",
        "🔍 Key Insights:",
        "  • Fewer incidents but longer resolution times",
        "  • Suggests more complex incident mix in October"
    ]
)

# Slide 9: What-If Analysis - Cumulative Impact
print("📊 Creating Slide 9: What-If Cumulative Impact...")
add_image_slide(
    "What-If Analysis: Cumulative Impact",
    os.path.join(OCT_PATH, "WhatIf_Cumulative_Impact.png"),
    "Top 5 events represent 36.3% of P75 TTM | Top 10 events represent 46.1%"
)

# Slide 10: What-If Analysis - Marginal Returns
print("📊 Creating Slide 10: What-If Marginal Returns...")
add_image_slide(
    "What-If Analysis: Diminishing Returns",
    os.path.join(OCT_PATH, "WhatIf_Cumulative_Marginal.png"),
    "Events 1-5: 13.8 min/event avg | Events 6-10: 3.7 min/event avg (73% less effective)"
)

# Slide 11: What-If Key Findings
print("📊 Creating Slide 11: What-If Key Findings...")
add_content_slide(
    "What-If Analysis: Key Findings",
    [
        "🎯 Event System Impact:",
        "  • 89 unique event systems identified",
        "  • 76 root events + 41 cascading outages",
        "",
        "📈 Top Event Impact:",
        "  • #1 Event: 5.3% of P75 TTM (10.0 minutes)",
        "  • Top 5 Events: 36.3% of P75 TTM (68.9 minutes)",
        "  • Top 10 Events: 46.1% of P75 TTM (87.6 minutes)",
        "",
        "⚡ Diminishing Returns:",
        "  • Events 1-5: Average 13.8 min reduction per event",
        "  • Events 6-10: Average 3.7 min reduction per event",
        "  • 73% reduction in effectiveness beyond top 5"
    ]
)

# Slide 12: Exclusions
print("📊 Creating Slide 12: Exclusions...")
add_content_slide(
    "Excluded Incidents",
    [
        "🚫 Exclusion Criteria:",
        "  • BCDR related incidents (planned drills)",
        "  • EUAP region incidents (pre-production)",
        "",
        "📊 Exclusion Summary:",
        "  • 3 incidents excluded (2.6% of total)",
        "  • 7,802 minutes TTM excluded",
        "  • Average TTM: 293 → 225 min (23% reduction)",
        "",
        "📋 Excluded Incidents:",
        "  • #694602140: Xstore, 6,095 min (EUAP)",
        "  • #694752515: Compute RP, 1,372 min (BCDR + EUAP)",
        "  • #694624704: SQL MI, 335 min (EUAP)"
    ]
)

# Slide 13: Narrative Insights - Resolution Gap
print("📊 Creating Slide 13: Narrative Insights...")
add_two_column_slide(
    "Narrative Insights: Resolution Method Gap",
    [
        "🔍 High TTM (≥P75):",
        "• Ad-hoc/Manual: 46.7%",
        "• TSG/Runbook: 10.0%",
        "• Automation: 0.0%",
        "• Avg TTM: 976 min",
        "",
        "Key Pattern:",
        "High TTM incidents rely heavily on manual investigation without documented procedures"
    ],
    [
        "✅ Normal TTM (<P75):",
        "• Ad-hoc/Manual: 20.7%",
        "• TSG/Runbook: 14.9%",
        "• Automation: 0.0%",
        "• Avg TTM: 58 min",
        "",
        "📊 The Gap:",
        "• +26.0pp ad-hoc usage",
        "• 17.0x longer resolution"
    ]
)

# Slide 14: Narrative Insights - Service Patterns
print("📊 Creating Slide 14: Service Patterns...")
add_content_slide(
    "Service-Specific Patterns",
    [
        "🔝 Top 3 Services by Total TTM Impact:",
        "",
        "1️⃣ SQL Control Plane: 12,140 min (41.4%)",
        "  • Service Fabric Failover Manager crashes",
        "  • No TSG coverage for assert failures",
        "  • Average: 3,035 minutes per incident",
        "",
        "2️⃣ Xstore: 7,322 min (25.0%)",
        "  • Storage connection timeouts",
        "  • Memory throttling issues",
        "  • Average: 1,220 minutes per incident",
        "",
        "3️⃣ SQL MI Prod Clusters: 1,793 min (6.1%)",
        "  • Certificate issues post-failover",
        "  • Average: 896 minutes per incident"
    ]
)

# Slide 15: Narrative Insights - Root Causes
print("📊 Creating Slide 15: Root Cause Patterns...")
add_two_column_slide(
    "Root Cause Analysis",
    [
        "High TTM Root Causes:",
        "• Network: 33.3%",
        "• Capacity/Resource: 30.0%",
        "• Hardware: 23.3%",
        "• Code/Software Bug: 20.0%",
        "• Deployment/Config: 6.7%",
        "",
        "⚠️ Critical Gap:",
        "Software bugs 5.9x more common in High TTM (20% vs 3.4%)"
    ],
    [
        "Detection Patterns:",
        "• BRAIN Detection: 0% (High TTM)",
        "• Customer-Reported: 6.7%",
        "• Internal-Monitoring: 10.0%",
        "",
        "🔍 Detection Gap:",
        "High TTM incidents bypass automated detection, indicating novel/complex failures"
    ]
)

# Slide 16: Recommendations
print("📊 Creating Slide 16: Recommendations...")
add_content_slide(
    "Actionable Recommendations",
    [
        "1️⃣ Create TSGs for Top Services (Priority 1)",
        "  • SQL Control Plane: Service Fabric troubleshooting",
        "  • Xstore: Storage timeout investigations",
        "  • SQL MI: Certificate validation post-failover",
        "",
        "2️⃣ Improve Detection Coverage (Priority 2)",
        "  • Add BRAIN monitors for capacity saturation",
        "  • Service Fabric assert failure detection",
        "  • 0% coverage in High TTM → opportunity area",
        "",
        "3️⃣ Address Code Quality (Priority 3)",
        "  • Focus on config validation logic",
        "  • Assert failure prevention",
        "  • Bugs 5.9x more common in High TTM",
        "",
        "4️⃣ Reduce Ad-hoc Resolution (Priority 4)",
        "  • Automate: certificate cleanup, rate limiting",
        "  • 46.7% ad-hoc rate → 26.0pp gap vs Normal TTM"
    ]
)

# Slide 17: Key Takeaways
print("📊 Creating Slide 17: Key Takeaways...")
add_content_slide(
    "Key Takeaways",
    [
        "✅ Volume Down, Complexity Up:",
        "  • 33% fewer incidents than September",
        "  • But P75 TTM increased 5.6% (more complex incidents)",
        "",
        "🎯 Focus on Top 5 Events:",
        "  • Represent 36.3% of total P75 TTM impact",
        "  • Diminishing returns beyond top 5",
        "",
        "🔧 Resolution Method is Key:",
        "  • High TTM: 46.7% ad-hoc (no TSG)",
        "  • Normal TTM: 20.7% ad-hoc",
        "  • Gap = 17.0x longer resolution time",
        "",
        "📊 Investment Priorities:",
        "  • SQL Control Plane, Xstore, SQL MI (67% of High TTM)",
        "  • TSG development > automation > detection"
    ]
)

# Slide 18: Appendix Title
print("📊 Creating Slide 18: Appendix...")
add_title_slide(
    "Appendix",
    "Detailed Analysis & Methodology"
)

# Slide 19: Data Sources
print("📊 Creating Slide 19: Data Sources...")
add_content_slide(
    "Data Sources & Methodology",
    [
        "📊 Data Source:",
        "  • Kusto Cluster: icmdataro.centralus.kusto.windows.net",
        "  • Database: IcmDataCommon",
        "  • Date Range: October 1-31, 2025",
        "  • Total Records: 117 incidents (487 columns)",
        "",
        "🔍 Analysis Methodology:",
        "  • Quantitative: Metrics, percentiles, What-If scenarios",
        "  • Qualitative: Text mining of incident narratives",
        "  • Event System Model: RootResponsibleIncidentId grouping",
        "",
        "🚫 Exclusions Applied:",
        "  • BCDR drills (1 incident)",
        "  • EUAP region (3 incidents)",
        "  • Filtered dataset: 114 incidents for analysis"
    ]
)

# Slide 20: Definitions
print("📊 Creating Slide 20: Definitions...")
add_content_slide(
    "Key Definitions",
    [
        "⏱️ TTM (Time to Mitigate):",
        "  • Total time from incident creation to mitigation",
        "",
        "📊 P75 (75th Percentile):",
        "  • 75% of incidents have TTM ≤ this value",
        "  • Key metric for measuring typical high-impact scenarios",
        "",
        "🔗 Event System:",
        "  • Root event + all cascading outages",
        "  • Grouped by RootResponsibleIncidentId",
        "",
        "🎯 High TTM Cohort:",
        "  • Incidents with TTM ≥ P75 (190 minutes)",
        "  • Used for comparative analysis vs Normal TTM",
        "",
        "🔧 Resolution Methods:",
        "  • Ad-hoc: Manual investigation without documented procedures",
        "  • TSG: Following documented troubleshooting guides",
        "  • Automation: Automated remediation without human intervention"
    ]
)

# Save presentation
print(f"\n💾 Saving presentation to: {OUTPUT_FILE}")
prs.save(OUTPUT_FILE)

print(f"\n{'='*80}")
print(f"✅ SUCCESS! PowerPoint presentation created with 20 slides")
print(f"📁 Location: {OUTPUT_FILE}")
print(f"{'='*80}\n")

print("📊 Slide Breakdown:")
print("  1. Title Slide")
print("  2. Executive Summary")
print("  3-7. Visualizations (TTM Distribution, Services, Timeline, Severity)")
print("  8. Month-over-Month Comparison")
print("  9-11. What-If Analysis (3 slides)")
print("  12. Exclusions")
print("  13-15. Narrative Insights (3 slides)")
print("  16. Recommendations")
print("  17. Key Takeaways")
print("  18-20. Appendix (3 slides)")
print(f"\n{'='*80}\n")
